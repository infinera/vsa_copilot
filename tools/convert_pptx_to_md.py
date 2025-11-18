#!/usr/bin/env python3
"""
Convert a PPTX file to Markdown with extracted images and notes.

Usage:
  python tools/convert_pptx_to_md.py input.pptx --out output.md --images-dir images --ocr

This script uses `python-pptx` to extract text and images. If `--ocr` is passed and
`pytesseract` + `Pillow` are available, it will attempt OCR on images and include
the recognized text in the Markdown.
"""
import argparse
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

try:
    from PIL import Image
except Exception:
    Image = None

try:
    import pytesseract
except Exception:
    pytesseract = None

import shutil
import subprocess
import sys

def ensure_dir(path):
    if not os.path.exists(path):
        os.makedirs(path, exist_ok=True)

def save_image(image, path):
    with open(path, 'wb') as f:
        f.write(image.blob)

def check_tesseract_available():
    return pytesseract is not None and shutil.which('tesseract') is not None

def extract_slide_text(slide):
    texts = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = []
        for paragraph in shape.text_frame.paragraphs:
            runs = [r.text for r in paragraph.runs]
            txt.append(''.join(runs))
        if txt:
            texts.append('\n'.join(txt))
    return '\n\n'.join(texts)

def extract_notes(slide):
    try:
        notes = slide.notes_slide.notes_text_frame.text
        return notes.strip()
    except Exception:
        return ''

def slide_title(slide):
    try:
        title_shape = slide.shapes.title
        if title_shape and title_shape.has_text_frame:
            return title_shape.text.strip()
    except Exception:
        pass
    return ''

def main():
    parser = argparse.ArgumentParser(description='Convert PPTX to Markdown')
    parser.add_argument('input', help='Path to input PPTX')
    parser.add_argument('--out', '-o', default='output.md', help='Output Markdown file')
    parser.add_argument('--images-dir', default='images', help='Directory to write extracted images')
    parser.add_argument('--ocr', action='store_true', help='Attempt OCR on extracted images (requires pytesseract+Pillow)')
    parser.add_argument('--slide-prefix', default='slide', help='Prefix for image filenames')
    args = parser.parse_args()

    prs = Presentation(args.input)
    ensure_dir(os.path.dirname(os.path.abspath(args.out)) or '.')
    ensure_dir(args.images_dir)

    md_lines = []
    md_lines.append('<!-- Converted from: {} -->\n'.format(os.path.basename(args.input)))

    for i, slide in enumerate(prs.slides, start=1):
        title = slide_title(slide)
        header = title if title else f'Slide {i}'
        md_lines.append(f'## {header}\n')

        body_text = extract_slide_text(slide)
        if body_text:
            md_lines.append(body_text + '\n')

        # extract images
        img_count = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = shape.image
                img_ext = img.ext
                img_count += 1
                img_name = f"{args.slide_prefix}-{i}-{img_count}.{img_ext}"
                img_path = os.path.join(args.images_dir, img_name)
                save_image(img, img_path)
                rel_path = os.path.relpath(img_path, os.path.dirname(os.path.abspath(args.out)) or '.')
                md_lines.append(f'![{img_name}]({rel_path})\n')
                # OCR if requested
                if args.ocr:
                    if not Image or not pytesseract:
                        md_lines.append('> OCR skipped: Pillow or pytesseract not installed\n')
                    elif not check_tesseract_available():
                        md_lines.append('> OCR skipped: tesseract executable not found in PATH\n')
                    else:
                        try:
                            ocr_text = pytesseract.image_to_string(Image.open(img_path))
                            if ocr_text.strip():
                                md_lines.append('> OCR: ' + ocr_text.strip().replace('\n', ' ') + '\n')
                        except Exception as e:
                            md_lines.append(f'> OCR failed: {e}\n')

        notes = extract_notes(slide)
        if notes:
            md_lines.append('\n**Notes:**\n')
            md_lines.append(notes + '\n')

        md_lines.append('\n---\n')

    with open(args.out, 'w', encoding='utf-8') as f:
        f.write('\n'.join(md_lines))

    print(f'Wrote Markdown to {args.out} and images to {args.images_dir}')

if __name__ == '__main__':
    main()
